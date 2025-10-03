import streamlit as st
import google.generativeai as genai
import os
from gtts import gTTS
import tempfile
import docx
from pptx import Presentation
import PyPDF2
from PIL import Image, ImageDraw, ImageFont
import textwrap
from moviepy import ImageClip, concatenate_videoclips, AudioFileClip
import datetime
import base64

# ----------------------------
# Configure Gemini API
# ----------------------------
api_key = os.environ.get("GEMINI_API_KEY")
genai.configure(api_key=api_key)

# ----------------------------
# File upload size limit
# ----------------------------
CUSTOM_MAX_SIZE_MB = 1024  # 1 GB

# ----------------------------
# Language codes
# ----------------------------
lang_codes = {
    "English": "en", "हिंदी": "hi", "Español": "es", "Français": "fr", "Deutsch": "de",
    "中文": "zh", "日本語": "ja", "한국어": "ko", "Русский": "ru", "اردو": "ur",
    "বাংলা": "bn", "తెలుగు": "te", "தமிழ்": "ta", "ગુજરાતી": "gu", "मराठी": "mr",
    "ਪੰਜਾਬੀ": "pa", "ಕನ್ನಡ": "kn", "मारवाड़ी": "rwr", "O‘zbekcha": "uz", "ქართული": "ka",
    "العربية": "ar", "Türkçe": "tr", "ภาษาไทย": "th", "فارسی": "fa", "Shqip": "sq",
    "Nederlands": "nl", "Svenska": "sv", "Italiano": "it", "Việt": "vi", "ລາວ": "lo"
}

# ----------------------------
# UI translations
# ----------------------------
def get_ui_texts(lang_code):
    translations = {
        "en": {"title":"TransformAI","paste_text":"Paste your text here","extra_comments":"Add extra instructions or comments (optional)",
               "upload_file":"Or upload a file (DOCX / PPTX / PDF / Audio)","output_type":"Choose output type:",
               "generate":"Generate","listen":"🔊 Listen","video":"🎥 Create Video","speech_lang":"Choose speech language:",
               "warning":"Please enter or upload some content.","output":"Output","try_speaker":"Try clicking the speaker below",
               "video_note":"📌 Note: Videos are text-to-slide only; users will no longer get 'I can’t generate videos' messages."},
        "hi": {"title":"ट्रांसफॉर्मAI","paste_text":"यहाँ अपना पाठ चिपकाएँ","extra_comments":"अतिरिक्त निर्देश या टिप्पणियाँ जोड़ें (वैकल्पिक)",
               "upload_file":"या एक फ़ाइल अपलोड करें (DOCX / PPTX / PDF / ऑडियो)","output_type":"आउटपुट प्रकार चुनें:",
               "generate":"जेनरेट करें","listen":"🔊 सुनें","video":"🎥 वीडियो बनाएँ","speech_lang":"वाक् भाषा चुनें:",
               "warning":"कृपया कुछ सामग्री दर्ज करें या अपलोड करें।","output":"आउटपुट","try_speaker":"नीचे दिए गए स्पीकर पर क्लिक करें",
               "video_note":"📌 नोट: वीडियो केवल टेक्स्ट-स्लाइड हैं; उपयोगकर्ता अब 'मैं वीडियो नहीं बना सकता' संदेश नहीं देखेंगे।"},
        "es": {"title":"TransformAI","paste_text":"Pega tu texto aquí","extra_comments":"Agregue instrucciones o comentarios adicionales (opcional)",
               "upload_file":"O sube un archivo (DOCX / PPTX / PDF / Audio)","output_type":"Elige tipo de salida:",
               "generate":"Generar","listen":"🔊 Escuchar","video":"🎥 Crear Video","speech_lang":"Elige idioma de voz:",
               "warning":"Por favor ingrese o cargue contenido.","output":"Salida","try_speaker":"Haz clic en el altavoz a continuación",
               "video_note":"📌 Nota: Los videos son solo diapositivas de texto; los usuarios ya no verán mensajes de 'No puedo generar videos'."},
    }
    return translations.get(lang_code, translations["en"])

# ----------------------------
# File extractors
# ----------------------------
def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def extract_text_from_audio(file_path):
    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
    try:
        return recognizer.recognize_google(audio_data)
    except:
        return "Could not understand the audio."

# ----------------------------
# Video generator
# ----------------------------
def create_text_video_with_audio(text, audio_lang="en", output_file="output_video.mp4"):
    chunks = textwrap.wrap(text, 200)
    clips = []
    temp_image_paths = []
    for i, chunk in enumerate(chunks):
        img = Image.new("RGB",(720,480),"white")
        draw = ImageDraw.Draw(img)
        font = ImageFont.load_default()
        y_text = 50
        for line in textwrap.wrap(chunk, width=40):
            draw.text((50,y_text), line, fill="black", font=font)
            y_text += 20
        tmp_img_path = os.path.join(tempfile.gettempdir(), f"slide_{i}.png")
        img.save(tmp_img_path)
        temp_image_paths.append(tmp_img_path)
        clips.append(ImageClip(tmp_img_path).set_duration(3))
    video = concatenate_videoclips(clips, method="compose")
    tts = gTTS(text=text, lang=audio_lang)
    tmp_audio = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
    tts.save(tmp_audio.name)
    video = video.set_audio(AudioFileClip(tmp_audio.name))
    video.write_videofile(output_file, fps=1)
    for path in temp_image_paths: os.remove(path)
    os.remove(tmp_audio.name)
    return output_file

# ----------------------------
# Session state
# ----------------------------
if "txt_input" not in st.session_state: st.session_state.txt_input = ""
if "extra_comments_input" not in st.session_state: st.session_state.extra_comments_input = ""
if "result" not in st.session_state: st.session_state.result = ""

# ----------------------------
# UI
# ----------------------------
selected_lang = st.selectbox("Choose language:", list(lang_codes.keys()))
selected_lang_code = lang_codes[selected_lang]
ui_texts = get_ui_texts(selected_lang_code)

st.title(ui_texts["title"])
txt_area = st.text_area(ui_texts["paste_text"], value=st.session_state.txt_input, height=150, key="txt_input")
extra_comments = st.text_area(ui_texts["extra_comments"], value=st.session_state.extra_comments_input, height=100, key="extra_comments_input")

# ----------------------------
# File upload
# ----------------------------
uploaded_file = st.file_uploader(ui_texts["upload_file"], type=["docx","pptx","pdf","mp3","wav"])
outpt = st.selectbox(ui_texts["output_type"], ["Summary","Quiz","Test","Audio","Video"])
if outpt=="Video": st.info(ui_texts["video_note"])

if uploaded_file:
    uploaded_file.seek(0)
    if uploaded_file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
        st.session_state.txt_input = extract_text_from_docx(uploaded_file)
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        st.session_state.txt_input = extract_text_from_pptx(uploaded_file)
    elif uploaded_file.type=="application/pdf":
        st.session_state.txt_input = extract_text_from_pdf(uploaded_file)
    elif uploaded_file.type in ["audio/mpeg","audio/wav"]:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            tmp.write(uploaded_file.read())
            tmp.flush()
            st.session_state.txt_input = extract_text_from_audio(tmp.name)

# ----------------------------
# Generate AI output
# ----------------------------
if st.button(ui_texts["generate"]):
    txt = st.session_state.txt_input
    extra_comments = st.session_state.extra_comments_input
    if not txt.strip():
        st.warning(ui_texts["warning"])
    else:
        prompt = f"Generate a {outpt} of the following text:\n\n{txt}"
        if extra_comments.strip(): prompt += f"\n\nInstructions: {extra_comments}"
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        st.session_state.result = response.candidates[0].content.parts[0].text
        st.subheader(ui_texts["output"])
        st.write(st.session_state.result)
